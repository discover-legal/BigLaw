"""Convert LAB task documents (.pdf/.docx/.xlsx/.pptx/text) to plain text.

The Go backend ingests text only (its /documents/upload returns 422 for PDF),
so all extraction happens here, client-side. Heavy libraries are imported
lazily so a missing optional dependency only breaks the formats that need it.
Falls back to `pandoc` on PATH for anything else, and skips (with a warning)
what it cannot read.
"""

from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

TEXT_EXTS = {".txt", ".md", ".markdown", ".csv", ".json", ".log", ".text", ".rtf", ".html", ".xml",
             ".eml", ".tsv", ".yaml", ".yml"}


class ConversionError(RuntimeError):
    pass


def _pdf_to_text(path: Path) -> str:
    try:
        import fitz  # PyMuPDF
    except ImportError as e:
        raise ConversionError("PyMuPDF not installed (pip install PyMuPDF)") from e
    with fitz.open(path) as doc:
        return "\n\n".join(page.get_text() for page in doc)


def _docx_to_text(path: Path) -> str:
    try:
        import docx
    except ImportError as e:
        raise ConversionError("python-docx not installed (pip install python-docx)") from e
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for table in d.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(parts)


def _xlsx_to_text(path: Path) -> str:
    try:
        import openpyxl
    except ImportError as e:
        raise ConversionError("openpyxl not installed (pip install openpyxl)") from e
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"## Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)


def _pptx_to_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ConversionError("python-pptx not installed (pip install python-pptx)") from e
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _pandoc_to_text(path: Path) -> str:
    if not shutil.which("pandoc"):
        raise ConversionError(f"no converter for {path.suffix} and pandoc not on PATH")
    proc = subprocess.run(
        ["pandoc", "--to=plain", str(path)],
        capture_output=True, text=True, timeout=120,
        # Force UTF-8: pandoc emits UTF-8, but text=True otherwise decodes with
        # the locale default (cp1252 on Windows), which dies on smart quotes etc.
        encoding="utf-8", errors="replace",
    )
    if proc.returncode != 0:
        raise ConversionError(f"pandoc failed on {path.name}: {proc.stderr.strip()[:200]}")
    return proc.stdout


def extract_text(path: Path) -> str:
    """Extract plain text from a task document. Raises ConversionError."""
    ext = path.suffix.lower()
    if ext in TEXT_EXTS:
        return path.read_text(encoding="utf-8", errors="replace")
    if ext == ".pdf":
        return _pdf_to_text(path)
    if ext == ".docx":
        return _docx_to_text(path)
    if ext in (".xlsx", ".xlsm"):
        return _xlsx_to_text(path)
    if ext == ".pptx":
        return _pptx_to_text(path)
    return _pandoc_to_text(path)


def convert_documents(documents_dir: Path) -> tuple[list[tuple[str, str]], list[str]]:
    """Walk a task's documents/ tree.

    Returns ([(relative_name, text), ...], [skipped_relative_names]).
    """
    converted: list[tuple[str, str]] = []
    skipped: list[str] = []
    if not documents_dir.is_dir():
        return converted, skipped

    for path in sorted(documents_dir.rglob("*")):
        if not path.is_file() or path.name.startswith("."):
            continue
        rel = str(path.relative_to(documents_dir))
        try:
            text = extract_text(path)
        except Exception as e:  # noqa: BLE001 — skip the file, never abort the run
            skipped.append(f"{rel} ({e})")
            continue
        if text and text.strip():
            converted.append((rel, text))
        else:
            skipped.append(f"{rel} (no extractable text)")
    return converted, skipped
